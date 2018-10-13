#!/usr/bin/env python

#
# File: pptx-creator.py
# Author: amort
#
# Dependencies:
#   - python-pptx: Install with "pip install --user python-pptx"
#   - openpyxl: Install with "pip install --user openpyxl"
#   - pathlib2: (python2 only) Install with "pip install --user pathlib2"
#
# Versions History:
#   Version 0.1 (7-30-2018)
#       Initial version.
#   Version 0.2 (9-12-2018)
#       Change input file to xml format
#       Add support for table generation -- TODO
#
# Description: This file takes an xml file with slide definitions and
#   creates a powerpoint presentation. The file can specify the format
#   of each slide, including slide title, text, images, etc. The
#   powerpoint is created using a template file, which can be specified
#   on the command line.
#
# Example: ./pptx-creator.py -t test/templates/blank -o test/example.pptx test/example.xml
#
# TODO
#   - work on text spacing between different values
#   - add table placeholder type
#   - add bulleted list placeholder type
#

import openpyxl # Import data from xlsx files

# Force python XML parser not faster C accelerators
# because we can't hook the C implementation
import sys
sys.modules['_elementtree'] = None

from pptx import Presentation
import xml.etree.ElementTree as ET
import datetime
import os
import re
import argparse

import csv      # Import data from csv files

try:
    import pathlib
except ImportError:
    import pathlib2 as pathlib

def main():
    # Parse arguments to get paths
    path_input, path_output, path_xml, path_pptx = parse_arguments()

    # Interpret template xml file
    template = get_template(path_xml)

    # Create presentation
    pc = PresentationCreator(path_pptx, template)
    pc.create_presentation(path_input, path_output)

# Parse input arguments
def parse_arguments():
    global verbose

    # Setup argument parser
    parser = argparse.ArgumentParser(description="" \
        "This program creates a powerpoint (pptx) file based on an xml"      \
        " definition file and a template. A template consists of a template" \
        " pptx file and a template xml file, both of which typically reside" \
        " in a template directory.")
    parser.add_argument("input",  help="xml definition file")
    parser.add_argument("-o", "--output", help="pptx output file")
    parser.add_argument("-t", "--template",
        help="template directory, containing pptx and xml files, "\
        "file names should either be template.pptx and template.xml or "\
        "<dirname>.pptx and <dirname>.xml"\
        "(specify individual file locations  with --pptx and --xml)")
    parser.add_argument("-p", "--pptx",
        help="template pptx file (overrides location of pptx from --template)")
    parser.add_argument("-x", "--xml",
        help="template xml file (overrides location of xml from --template)")
    parser.add_argument("-v", "--verbose", action="store_true")
    args = parser.parse_args()

    # Get arguments
    verbose           = args.verbose
    filename_input    = args.input
    filename_output   = args.output
    filename_template = args.template
    filename_pptx     = args.pptx
    filename_xml      = args.xml

    # Get pathlib object for input
    path_input = pathlib.Path(filename_input)

    if verbose:
        print("INFO: Input path " + str(path_input) + ".")

    # Get pathlib object for output
    if (filename_output):
        path_output = pathlib.Path(filename_output)
    else:
        # Path to output file with same name as input file in pwd
        path_output = path_input.with_suffix('.pptx').name

    if verbose:
        print("INFO: Output path " + str(path_output) + ".")

    # Get pathlib object for template
    if (filename_template):
        path_template = pathlib.Path(filename_template)
        template_names = (path_template.name, 'template')

    # Get xml path
    if (filename_xml):
        path_xml = pathlib.Path(filename_xml)
    else:
        if (not path_template):
            raise ValueError("No xml or template path specified!")

        # Search for valid xml paths
        for name in template_names:
            path_temp = path_template.joinpath(name + '.xml')
            if path_temp.exists():
                path_xml = path_temp
                break

        # Unable to find valid xml file
        if (not path_xml):
            raise ValueError("No valid template xml file found!")

    if verbose:
        print("INFO: XML template path " + str(path_xml) + ".")

    # Get pptx path
    if (filename_pptx):
        path_pptx = pathlib.Path(filename_pptx)
    else:
        if (not path_template):
            raise ValueError("No pptx or template path specified!")

        # Search for valid xml paths
        for name in template_names:
            path_temp = path_template.joinpath(name + '.pptx')
            if path_temp.exists():
                path_pptx = path_temp
                break

        # Unable to find valid xml file
        if (not path_pptx):
            raise ValueError("No valid template pptx file found!")

    if verbose:
        print("INFO: PPTX template path " + str(path_pptx) + ".")

    return path_input, path_output, path_xml, path_pptx

# Get template mapping from xml file
def get_template(path_xml):
    template = {}

    # Open xml template definition file
    xml_template = ET.parse(str(path_xml))

    # Create map from child (c) to parent (p)
    xml_parents = {c:p for p in xml_template.iter() for c in p}

    template_idx = -1
    layout_idx = -1
    ph_idx = -1

    # Iterate over xml elements to create template
    for xml_elem in xml_template.iter():
        if (xml_elem.tag == 'template'):
            template_idx += 1

            # Ensure template has no parent
            try:
                xml_parents[xml_elem]
            except KeyError:
                pass
            else:
                raise ValueError("Template should be top level element!")

            if (template_idx > 0):
                raise ValueError("Cannot have more than one template element!")

        elif (xml_elem.tag == 'layout'):
            # Update index values for error messages
            layout_idx += 1
            ph_idx = -1

            layout = xml_elem.get("name")
            index  = xml_elem.get("index")

            if (not layout):
                raise ValueError("Layout \"{}\" missing name attribute in xml "\
                    "template file {}"\
                    "".format(layout_idx, str(path_xml)))

            if (xml_parents[xml_elem].tag != 'template'):
                raise ValueError("Layout \"{}\" ({}) has parent node "\
                    "\"{}\" and should have parent node \"template\" in xml "\
                    " template file {}"\
                    "".format(layout_idx, layout, xml_parents[xml_elem].tag,
                        str(path_xml)))

            if (not index):
                raise ValueError("Layout \"{}\" ({}) missing index attribute "\
                    "in xml template file {}"\
                    "".format(layout_idx, layout, str(path_xml)))

            # Associate index with layout name
            template[layout] = {}
            template[layout]["idx"] = int(index)
            template[layout]["ph"] = {}

        elif (xml_elem.tag == 'placeholder'):
            # Update index values for error messages
            ph_idx += 1

            ph     = xml_elem.get("name")
            index  = xml_elem.get("index")

            if (not ph):
                raise ValueError("Placeholder \"{}\" in layout \"{}\" ({}) "\
                    "missing name attribute in xml template file {}"\
                    "".format(ph_idx, layout_idx, layout, str(path_xml)))

            if (xml_parents[xml_elem].tag != 'layout'):
                raise ValueError("Placeholder \"{}\" ({}) has parent node "\
                    "\"{}\" and should have parent node \"layout\" in xml "\
                    " template file {}"\
                    "".format(ph_idx, ph, xml_parents[xml_elem].tag,
                        str(path_xml)))

            if (not index):
                raise ValueError("Placeholder \"{}\" ({}) in layout \"{}\" "\
                    "({}) missing index attribute in xml template file {}"\
                    "".format(ph_idx, ph, layout_idx, layout, str(path_xml)))

            # Associate index with ph name for current layout
            template[layout]["ph"][ph] = int(index)

        else:
            raise ValueError("Invalid tag \"{}\" in xml template file {}"\
                "".format(xml_elem.tag, str(path_xml)))

    # Invalid template file if no layout tags are found
    if (layout_idx < 0):
        raise ValueError("No layout tags found in xml template file {}"\
            "".format(str(path_xml)))

    return template

class PresentationCreator:
    """ The presentation creator class is used to create a pptx presentation.

    This class is used to create a pptx presentation for the pptx-creator
    project. It utilizes the pptx module to modify a template pptx file
    and create a presentation as specified in an input XML file.

    """

    ph_types = {"text", "image", "table"}


    def __init__(self, path_pptx, template):
        self.path_pptx   = path_pptx
        self.template    = template

    def create_presentation(self, path_input, path_output):
        """ Create pptx presentation

        Create a new presentation based on the input from the path_input XML
        file and save it to path_output.

        Args:
            path_input: path to input XML file
            path_output: path to save pptx file
        """

        self.invalid_images = []
        self.invalid_imports = {}

        self.input  = path_input
        self.output = path_output

        # Create presentation
        self.prs = Presentation(str(self.path_pptx))

        # Process the input xml file
        self.ppp = PresentationPreprocessor(str(self.input))

        # Create slides and define slide references
        self._initialize_slides(self.prs, self.ppp.get_root())

        # Process slides and fill fields
        self._process_slides(self.prs, self.ppp.get_root())

        # Report missing image paths
        if (len(self.invalid_images) > 0):
            print("\nInvalid Image Paths: ")
            for img in self.invalid_images:
                print("  " + img)

        # Report invalid table entries
        if (len(self.invalid_imports) > 0):
            print("\nInvalid Import: ")
            for path,errors in self.invalid_imports.items():
                print("  " + path)
                for msg in errors:
                    print("    " + msg)

        self.prs.save(str(self.output))

        print("\nPresentation created: {}\n".format(self.output))

    def _initialize_slides(self, prs, root_entry):
        """ Create empty slides and save references

        Iterate through slide entries from the preprocessor and create
        empty slides with the specified layouts and create a dictionary
        of references for making links between slides. This must be
        done before processing slides because all slides to be referenced
        must be created before the link is added to a slide.

        Args:
            prs: Presentation object where slides will be created
            root_entry: root PreprocessorEntry object containing input tree

        """

        # Initialize data structures
        self.slides      = []
        self.layouts     = []
        self.refs        = {}

        # Iterate through slides
        for slide in root_entry.data:
            # Get slide layout
            layout_vals = slide.get_values(tag="layout", join=True)

            if (len(layout_vals) > 1):
                raise ValueError("slide may only have one layout attribute"\
                        "\n{}".format(self.ppp.error_info(slide)))

            # Add layout to list and remove layout entry from slide
            self.layouts.append(layout_vals[0])
            slide.remove(tag="layout")

            # Ceate new slide with layout
            prs_layout = prs.slide_layouts[self.template[layout_vals[0]]["idx"]]
            self.slides.append(prs.slides.add_slide(prs_layout))

            # Get slide reference label
            label_vals = slide.get_values(tag="label", join=True)

            # Reference label points to current slide and remove label entry
            if (len(label_vals) == 1):
                self.refs[label_vals[0]] = self.slides[-1]
                slide.remove(tag="label")
            elif (len(label_vals) > 1):
                raise ValueError("slide may only have one label attribute"\
                        "\n{}".format(self.ppp.error_info(slide)))


    def _process_slides(self, prs, root_entry):
        """ Add data to the slides that were previously created

        Iterate through slide entries from the preprocessor and add
        information to the slides created by the _initiliaze_slides
        subroutine. This must be done is two steps because slides to be
        referenced through links must be created before the links to
        that slide is created.

        Args:
            prs: Presentation object where slides will be created
            root_entry: root PreprocessorEntry object containing input tree

        """

        # Iterate through slides
        for i,slide in enumerate(root_entry.data):

            # Iterate through placeholders
            for ph in slide.data:
                # Ensure no value entries directly under slide
                if (ph.which != "entry"):
                    raise ValueError("slide may only contain placeholder"\
                            "elements"\
                            "\n{}".format(self.ppp.error_info(ph)))

                # Get placeholder from slide
                try:
                    ph_idx = self.template[self.layouts[i]]["ph"][ph.tag]
                except KeyError:
                    raise ValueError("placeholder \"{}\" not found in "\
                            "template\n{}"\
                            "".format(ph.tag, self.ppp.error_info(ph)))

                prs_ph = self.slides[i].placeholders[ph_idx]

                # Determine if placeholder has a type
                type_vals = ph.get_values(tag="type", join=True)

                if (len(type_vals) > 1):
                    raise ValueError("placeholder may only have one type "\
                            "attribute\n{}"\
                            "".format(self.ppp.error_info(ph)))
                else:
                    # No type specified, find type elements under placeholder
                    if (len(type_vals) == 0):
                        found = 0
                        for sub in ph.data:
                            if (sub.which == "value"):
                                continue

                            # Determine if entry is a type element
                            for t in self.ph_types:
                                if (sub.tag == t):
                                    ph = sub
                                    type_vals.insert(0,t)
                                    found += 1

                        # No valid type found, assume data is text
                        if (found == 0):
                            type_vals.insert(0,"text")

                        # Found more than one type, this is invalid
                        elif (found > 1):
                            raise ValueError("placeholder may only have one "\
                                    "type\n{}".format(self.ppp.error_info(ph)))

                    # Type was specified as an attribute of the placeholder
                    else:
                        # remove type entry from placeholder entry
                        ph.remove(tag="type")

                    # Ensure class has type function
                    if (not hasattr(self, "_ph_" + type_vals[0])):
                        raise ValueError("placeholder type \"{}\" is not "\
                                "valid\n{}".format(type_vals[0],
                                    self.ppp.error_info(ph)))

                    # Call type function
                    type_func = getattr(self, "_ph_" + type_vals[0])
                    type_func(ph, prs_ph, self.slides[i])

    def _prs_insert_text(self, entry, prs_object):
        """ Add text to a shape

        Process text contained in the PreprocessorEntry and place it in the
        prs_object containing a text_frame where the text should be inserted.

        Args:
            entry: PreprocessorEntry with text data
            prs_object: presentation object with a text_frame where text is
                to be inserted

        """

        text = ""

        for sub in entry.data:
            if (sub.which == "value"):
                text += str(sub.value)
                continue

            if (sub.tag == "date"):
                text += datetime.datetime.now().strftime("%B %d, %Y")
            else:
                raise ValueError("invalid \"{}\" entry in text placeholder."\
                        "\n{}".format(sub.tag, self.ppp.error_info(sub)))

        prs_object.text = text

    def _ph_text(self, entry, prs_ph, prs_slide):
        """ Add text to the placeholder

        Process text contained in the PreprocessorEntry and place it in the
        prs_ph placeholder on the prs_slide slide.

        Args:
            entry: PreprocessorEntry with text data
            prs_ph: presentation placeholder where text is to be inserted
            prs_slide: presentation slide containing placeholder

        """

        self._prs_insert_text(entry, prs_ph)

    def _ph_image(self, entry, prs_ph, prs_slide):
        """ Add image to the placeholder

        Process image path contained in the PreprocessorEntry and place image
        on the prs_slide slide at the location of the prs_ph placeholder.

        The general purpose placeholder cannot hold a picture in the current
        version of the pptx module so a new picture shape is added at the
        position of the placeholder and the original placeholder is deleted.

        If the image path specified in input XML file is invlalid, a not
        indicating this is added to the placeholder and the image is added
        to an invalid image list.

        Args:
            entry: PreprocessorEntry with text data
            prs_ph: presentation placeholder where text is to be inserted
            prs_slide: presentation slide containing placeholder

        """

        path = ""

        # Create path
        for sub in entry.data:
            if (sub.which == "value"):
                path += sub.value
                continue

            # No valid sub elements in image tag
            raise ValueError("invalid \"{}\" entry in image placeholder."\
                    "\n{}".format(sub.tag, self.ppp.error_info(sub)))

        if not pathlib.Path(path).exists():
            self.invalid_images.append(path)
            prs_ph.text = "Image Not Found: " + path
            return

        # add picture in a new picture shape at location of placeholder
        pic = prs_slide.shapes.add_picture(path, prs_ph.left, prs_ph.top)

        # calculate size to fit inside placeholder area
        ratio = min(prs_ph.width  / float(pic.width), prs_ph.height / float(pic.height))

        pic.height = int(pic.height * ratio)
        pic.width  = int(pic.width  * ratio)

        pic.left = int(prs_ph.left + ((prs_ph.width  - pic.width)/2))
        pic.top  = int(prs_ph.top  + ((prs_ph.height - pic.height)/2))

        # remove placeholder from slide
        elem = prs_ph.element
        elem.getparent().remove(elem)

    def _ph_table(self, entry, prs_ph, prs_slide):
        """ Add table to the placeholder

        Process the table description contained in the PreprocessorEntry
        and place table on the prs_slide slide at the location of the
        prs_ph placeholder.

        The general purpose placeholder cannot hold a table in the current
        version of the pptx module so a new table shape is added at the
        position of the placeholder and the original placeholder is deleted.

        Args:
            entry: PreprocessorEntry with text data
            prs_ph: presentation placeholder where text is to be inserted
            prs_slide: presentation slide containing placeholder

        """

        table = []
        max_col = 0
        cur_row = -1

        # Iterate through rows
        for row in entry.data:
            if (row.which == "value"):
                raise ValueError("invalid value \"{}\" in table "\
                        "placeholder, expected row element.\n{}"\
                        "".format(row.value, self.ppp.error_info(row)))

            if (row.tag == "row"):
                # Initialize new row in array
                cur_row += 1
                cur_col = -1
                if len(table) <= cur_row:
                    table.insert(cur_row, [])

                # Iterate through columns
                for col in row.data:
                    if (col.which == "value"):
                        raise ValueError("invalid value \"{}\" in table "\
                                "placeholder, expected cell element.\n{}"\
                                "".format(row.value, self.ppp.error_info(col)))

                    # Add cell entry to table array
                    if (col.tag == "cell"):
                        while(True):
                            cur_col += 1

                            # Add data if at end of row
                            if len(table[cur_row]) <= cur_col:
                                table[cur_row].append(col)
                                break

                            # Add data if current cell is empty
                            elif table[cur_row][cur_col] is None:
                                table[cur_row][cur_col] = col
                                break

                        # TODO - delete me
                        #print ("row={}\n  {}".format(cur_row, str(table[cur_row])))

                    # Add import entries to table array
                    elif (col.tag == "import"):
                        cur_col += 1

                        # TODO -- need to fix tables spanning multiple rows
                        # Add entries to table without changing cur_row
                        for i,irow in enumerate(self._import(col)):
                            if len(table) <= cur_row + i:
                                table.insert(cur_row + i, [])

                            # Add preceding cells when cur_col != 0
                            row_len = len(table[cur_row + i])
                            if (row_len < cur_col):
                                for j in range(row_len,cur_col):
                                    table[cur_row + i].insert(j, None)

                            # Insert entries into table
                            for j,icol in enumerate(irow):
                                table[cur_row + i].insert(cur_col + j, icol)

                            # TODO - delete me
                            #print ("row={}\n  {}".format(cur_row + i, str(table[cur_row + i])))

                    # Invalid element
                    else:
                        raise ValueError("invalid element \"{}\" in table "\
                                "row, expected cell element.\n{}"\
                                "".format(col.tag, self.ppp.error_info(col)))

                # Update max_col count
                if (len(table[cur_row]) > max_col+1):
                    max_col = len(table[cur_row]) - 1

            # Add import entries to table array
            elif (row.tag == "import"):

                # Add entries to table and change cur_row
                for irow in self._import(row):
                    cur_row += 1
                    if len(table) <= cur_row:
                        table.insert(cur_row, [])

                    for icol in irow:
                        table[cur_row].append(icol)

            else:
                raise ValueError("invalid element \"{}\" in table, "\
                        "expected row element.\n{}"\
                        "".format(row.tag, self.ppp.error_info(row)))

        # Create table on slide at location of placeholder
        prs_table = prs_slide.shapes.add_table(len(table), max_col+1,
                prs_ph.left, prs_ph.top, prs_ph.width, prs_ph.height).table

        # remove placeholder from slide
        elem = prs_ph.element
        elem.getparent().remove(elem)

        # place text into table
        for i in range(0,len(table)):
            for j in range(0,max_col+1):
                # Some rows may not contain all columns
                if (j >= len(table[i])):
                    break

                # Some cells may be empty
                elif (table[i][j] is None):
                    continue

                #TODO - delete me
                #print("row:{},col:{},data:{}".format(i, j, table[i][j]))

                # Add text from element to cell
                self._prs_insert_text(table[i][j], prs_table.cell(i,j))

    def _import(self, entry):
        """ Import data from file for presentation

        Read the file based on the extension (.csv, .xlsx) and return an
        array of entries containing the data that is read. The name of the
        file is specified as the value of the import element.

        For each file category, attributes can be used to specify which
        data should be imported:

            Spreadsheet(.csv/.xlsx)
                row: may be a single value, range, or list to indicate which
                    rows should be imported
                col: may be a single value, range, or list to indicate which
                    columns should be imported

        Args:
            entry: PreprocessorEntry with text data

        Return:
            2 dimensional array of entries containing the imported data

        """

        cat = None

        # Get filename from entry
        path_file = pathlib.Path(entry.get_values(join=True))

        # Determine type of file and open appropriate importer
        if (path_file.suffix == ".xlsx"):
            cat = "spreadsheet"

            importer = ImportXLSX(str(path_file), entry=entry)

        elif (path_file.suffix == ".csv"):
            cat = "spreadsheet"

            importer = ImportCSV(str(path_file), entry=entry)

        else:
            raise ValueError("invalid import suffix \"{}\"\n{}"\
                    "".format(path_file.suffix, self.ppp.error_info(entry)))

        # Spreadsheet file
        if (cat == "spreadsheet"):
            num_row = None
            num_col = None

            for child in entry.data:
                # Ignore value entries (we already have filename
                if (child.which == "value"):
                    continue

                # Get number of rows expected
                elif (child.tag == "num_row"):
                    num_row = child.get_values(join=True)

                # Get number of columns expected
                elif (child.tag == "num_col"):
                    num_col = child.get_values(join=True)

                # Get row spec
                elif (child.tag == "row"):
                    importer.add_row(child.get_values(join=True))

                # Get col spec
                elif (child.tag == "col"):
                    importer.add_col(child.get_values(join=True))

                # Get row key
                elif (child.tag == "row_key"):
                    col = None
                    re  = False

                    # Get associated values
                    for attr in child.data:
                        if (attr.which == "value"):
                            continue
                        elif (attr.tag == "col"):
                            col = attr.get_values(join=True)
                        elif (attr.tag == "re"):
                            re = attr.get_values(join=True)
                        else:
                            raise ValueError("invalid row_key import "\
                                    "attribute \"{}\"\n{}"\
                                    "".format(attr.tag,
                                        self.ppp.error_info(attr)))

                    importer.add_row_key(child.get_values(join=True),
                        col=col, re=re)

                # Get col key
                elif (child.tag == "col_key"):
                    row = None
                    re  = False

                    # Get associated values
                    for attr in child.data:
                        if (attr.which == "value"):
                            continue
                        elif (attr.tag == "row"):
                            row = attr.get_values(join=True)
                        elif (attr.tag == "re"):
                            re = attr.get_values(join=True)
                        else:
                            raise ValueError("invalid col_key import "\
                                    "attribute \"{}\"\n{}"\
                                    "".format(attr.tag,
                                        self.ppp.error_info(attr)))

                    importer.add_col_key(child.get_values(join=True),
                        col=col, re=re)

                # Get sheet
                elif (child.tag == "sheet"):
                    importer.add_sheet(child.get_values(join=True))

                # Invalid child element
                else:
                    raise ValueError("invalid import attribute \"{}\"\n{}"\
                            "".format(child.tag, self.ppp.error_info(child)))


            # Get spreadsheet data importer
            error = False
            entries = [[]]
            try:
                entries = importer.read()
            except (IndexError, KeyError) as err:
                error = True

                # Add error message
                msg = "{} \"{}\": {}".format(err.args[1],
                        err.args[2], err.args[0])
            except ValueError as err:
                error = True

                # Add error message
                msg = "{}".format(err.args[0])

            # Add error message to dictionary
            if error:
                if not str(path_file) in self.invalid_imports:
                    self.invalid_imports[str(path_file)] = []

                self.invalid_imports[str(path_file)].append(msg)

            # Get number of rows from importer if not set
            if num_row is None:
                num_row = importer.num_row

                if num_row is None:
                    num_row = 1

            # Get number of columns from importer if not set
            if num_col is None:
                num_col = importer.num_col

                if num_col is None:
                    num_col = 1


            # Create temporary PreprocessEntry for any unspecified cells
            tmp = PreprocessorEntry(entry.tag, parent=entry,
                    elem=entry.elem, value="N/A")

            # Ensure size of entries matches number of columns and rows
            valid_entries = []
            for i in range(0, num_row):
                valid_entries.append([])

                for j in range(0, num_col):
                    if (len(entries) > i and len(entries[i]) > j):
                        valid_entries[i].append(entries[i][j])
                    else:
                        valid_entries[i].append(tmp)

            # TODO -- show warnings for num col/row mismatch
            # TODO -- show line numbers for invalid imports
            # Check for change in size of rows
            if (len(valid_entries) != len(entries)):
                msg = "row number mismatch: expected {}, got {}"\
                    "".format(len(valid_entries), len(entries))

                # TODO
                #if not str(path_file) in self.invalid_imports:
                    #self.invalid_imports[str(path_file)] = []

                #self.invalid_imports[str(path_file)].append(msg)

            # Check for change in size of columns
            entries_col = (0 if len(entries) == 0 else len(entries[0]))
            if (len(valid_entries[0]) != entries_col):
                msg = "column number mismatch: expected {}, got {}"\
                    "".format(len(valid_entries), entries_col)

                # TODO
                #if not str(path_file) in self.invalid_imports:
                    #self.invalid_imports[str(path_file)] = []

                #self.invalid_imports[str(path_file)].append(msg)

        # Invalid category
        else:
            raise AssertionError("invalid import category \"{}\""\
                    "".format(cat))

        return valid_entries

class ImportSpreadsheet(object):
    """ Inherited by classes uses to read specific Spreadsheet file formats

    Contains the base functions shared by all types of spreadsheets

    """

    def __init__(self, filename, filetype, entry=None):
        self.filename=filename
        self.filetype=filetype
        self.entry = entry
        self.data = None
        self.rows = []
        self.cols = []
        self.row_keys = []
        self.col_keys = []

    def add_row(self, row):
        """ Set the row or rows to be read from the spreadsheet

        A string representing the row(s) to be read from the spreadsheet
        is passed to this function and stored in this object until the data
        is to be read from the file. This function can be called multiple
        times to add to the number of rows that should be read. The data from
        the rows will be returned in the order they are specified and the data
        from a row can be read multiple times. If no columns are specified,
        all the columns in the spreadsheet will be returned.

        Rows in a file start at 1 and count up.

        Row specification examples:
            "1"     = read row 1 from the file
            "1:6"   = read rows 1 through 6 from the file
            "2,4,2" = read rows 2, 4, and 2 from the file
            "2,4-6" = read rows 2, 4, 5, 6 from the file

        Args:
            row: string specifying the rows to be read

        """

        # Add values from row spec to row array
        self.rows += self._get_array(row)

    def add_col(self, col):
        """ Set the column or columns to be read from the spreadsheet

        A string representing the column(s) to be read from the spreadsheet
        is passed to this function and stored in this object until the data
        is to be read from the file. This function can be called multiple
        times to add to the number of columns that should be read. The data
        from the columns will be retrieved in the order they are specified 
        and the data from a column can be read multiple times. If no
        columns are specified, all the columns in the spreadsheet will
        be returned.

        Columns in a file start at 1 (or A) and count up.

        Column specification examples:
            "a"     = read the first column from the file
            "a:6"   = read columns 1 through 6 from the file
            "2,d,7" = read columns 2, 4, and 7 from the file
            "2,4-6" = read columns 2, 4, 5, 6 from the file

        Args:
            col: string specifying the rows to be read

        """

        # Add values from column spec to column array
        self.cols += self._get_array(col)

    def add_row_key(self, row_key, col=None, re=False):
        """ Set the key used to match rows in the spreadsheet

        This adds a key to the list of keys used to filter the rows that
        are returned from the spreadsheet. The rows originally in the
        set to be searched are those that were requested by the add_row()
        function, if no rows are specified, all rows will be searched.
        By default, the filter will search through all the columns for
        each row looking for a match, regardless of the columns
        requested by the add_col() function. To limit the columns
        searched for each row, a col spec can be passed to this function
        and the filtering will only look at those specified columns.
        If a row matches the filtering, only the columns specified
        by the add_col() function will be returned.

        Args:
            row_key: a string used to select rows to return when that
                row has data that matches the string
                (by default all columns are searched unless col is specified)

        Kwargs:
            col: string specifying which columns to use when finding
                a match to the row_key
            re: indicates whether regular expression foramtting is used
                for the row_key string

        """

        self.row_keys.append({})
        self.row_keys[-1]["re"] = re

        # Create re object when requested
        if re:
            self.row_keys[-1]["key"] = re.compile(row_key)
        else:
            self.row_keys[-1]["key"] = row_key

        # Set col spec if specified
        if not col is None:
            self.row_keys[-1]["col"] = self._get_array(col)
        else:
            self.row_keys[-1]["col"] = None


    def add_col_key(self, col_key, row=None, re=False):
        """ Set the key used to match columns in the spreadsheet

        This adds a key to the list of keys used to filter the columns that
        are returned from the spreadsheet. The columns originally in the
        set to be searched are those that were requested by the add_col()
        function, if no columns are specified, all columns will be searched.
        By default, the filter will search through all the rows for
        each column looking for a match, regardless of the rows
        requested by the add_row() function. To limit the rows
        searched for each column, a row spec can be passed to this function
        and the filtering will only look at those specified rows.
        If a column matches the filtering, only the rows specified
        by the add_row() function will be returned.

        Args:
            col_key: a string used to select columns to return when that
                column has data that matches the string
                (by default all rows are searched unless row is specified)

        Kwargs:
            row: string specifying which rows to use when finding
                a match to the col_key
            re: indicates whether regular expression foramtting is used
                for the col_key string

        """

        self.col_keys.append({})
        self.col_keys[-1]["re"]  = re

        # Create re object when requested
        if re:
            self.col_keys[-1]["key"] = re.compile(col_key)
        else:
            self.col_keys[-1]["key"] = col_key

        # Set row spec if specified
        if not row is None:
            self.col_keys[-1]["row"] = self._get_array(row)
        else:
            self.col_keys[-1]["row"] = None

    def get_data(self):
        """ Get data specified by column and row specifications

        Return data as a 2-dimensional array of PreprocessorEntries as
        filtered by the column and row specifications.

        Return:
            2-dimensional array of data as PreprocessorEntries

        """

        # Data must be set
        if self.data is None:
            raise AssertionError("data must be set before calling get_data()")


        # Create rows/cols array if they haven't been set
        if len(self.rows) == 0:
            self.rows = list(range(1, len(self.data)+1))

        if len(self.cols) == 0:
            if len(self.data) > 0:
                self.cols = list(range(1, len(self.data[0])+1))
            else:
                self.cols = []


        # Save number of rows/columns for error messages
        if len(self.row_keys) == 0:
            self.num_row = 1
        else:
            self.num_row = len(self.rows)
            if self.num_row < 1:
                self.num_row = 1

        if len(self.col_keys) == 0:
            self.num_col = 1
        else:
            self.num_col = len(self.cols)
            if self.num_col < 1:
                self.num_col = 1


        # Data is empty, return empty array
        if (len(self.data) < 1 or len(self.data[0]) < 1):
            raise ValueError("no data found")


        # Go through each row_key and filter data
        for rk in self.row_keys:
            rk_match = False

            # Get columns to search
            rk_cols = rk["col"]
            if rk_cols is None:
                rk_cols = list(range(1, len(self.data[0])+1))

            # Only search specified rows
            valid_rows = list(self.rows)
            for r in self.rows:
                match=False

                # Search through each column for key
                for c in rk_cols:
                    val = self.data[r-1][c-1].get_values(join=True)

                    # Search for match
                    if (rk["re"]):
                        match = rk["key"].match(val)
                    else:
                        match = val.find(rk["key"]) > -1

                    # TODO - delete
                    #print("k:{}, r:{}, c:{}, match:{}".format(str(rk["key"]),
                    #    r, c, (1 if match else 0)))

                    if match:
                        break

                # Remove row if it doesn't match
                if not match:
                    valid_rows.remove(r)
                else:
                    rk_match = True

            # Update rows with valid_rows
            self.rows = valid_rows

            # When there
            if not rk_match:
                raise KeyError("no match found", "row key", str(rk["key"]))



        # Update number of rows
        self.num_row = len(self.rows)
        if self.num_row < 1:
            self.num_row = 1


        # Go through each col_key and filter data
        for ck in self.col_keys:
            ck_match = False

            # Get rows to search
            ck_rows = ck["row"]
            if ck_rows is None:
                ck_rows = list(range(1, len(self.data)+1))

            # Only search specified columns
            valid_cols = list(self.cols)
            for c in self.cols:
                match=False

                # Search through each row for key
                for r in ck_rows:
                    val = self.data[r-1][c-1].get_values(join=True)

                    # Search for match
                    if (ck["re"]):
                        match = ck["key"].match(val)
                    else:
                        match = val.find(ck["key"]) > -1

                    if match:
                        break

                # Remove column if it doesn't match
                if not match:
                    valid_cols.remove(c)
                else:
                    ck_match = True

            # Update cols with valid_cols
            self.cols = valid_cols

            # When there
            if not ck_match:
                raise KeyError("no match found", "column key", str(rk.key))


        # Update number of columns
        self.num_col = len(self.cols)
        if self.num_col < 1:
            self.num_col = 1


        # Generate return array
        ret = []
        for i,r in enumerate(self.rows):
            # Check if in range
            if r > len(self.data):
                raise IndexError("index out of range", "row", str(r))

            ret.append([])
            for c in self.cols:
                # Check if in range
                if c > len(self.data[r-1]):
                    raise IndexError("index out of range", "column", str(c))

                ret[i].append(self.data[r-1][c-1])

        return ret

    def _get_array(self, spec):
        """ Return array of values specified by the row or column spec

        The spec string contains information about which columns or rows
        should be read from a file. This function converts that spec into
        an array of values that can be used to read the data.

        Example specifications:
            "a"     = returns [1]
            "a:6"   = returns [1,2,3,4,5,6]
            "2,d,7" = returns [2,4,7]
            "2,4-6" = returns [2,4,5,6]
            "2:12:4" = returns [2,6,10]

        Args:
            spec: string specifying row/column ranges

        Return:
            array derivied from spec row/column ranges

        """

        array = []

        # Define re to split range specs
        re_range   = re.compile('\s*[-:]\s*')
        re_digits  = re.compile('[0-9]')
        re_letters = re.compile('[a-z]')

        # Split spec at commas (don't split range specs)
        # TODO fix this so it never matches ""
        list_vals = re.split('(?<![-:])\s*,?\s*(?![-:])', spec)

        # Loop through split spec values
        for list_val in list_vals:

            # Split all values into range specs if possible
            range_vals = re_range.split(list_val)

            # Convert all values from letters to numbers
            for i,range_val in enumerate(range_vals):
                tmp = range_val.lower()

                # All letters, convert to int
                if re_letters.match(tmp):
                    # Don't combine letters and digits in single value
                    if re_digits.match(tmp):
                        raise ValueError('cannot mix digits and letters "{}" '\
                            'in row/column spec "{}"'.format(range_val, spec))

                    num = 0
                    for j,char in enumerate(tmp):
                        tmp = 26**(len(tmp) - j - 1)
                        tmp *= ord(char) - ord('a') + 1
                        num += tmp

                    tmp = num

                # All digits, interpret at int
                else:
                    tmp = int(range_val)

                # Value must be greater 0
                if tmp < 1:
                    raise ValueError("index \"{}\" cannot be less than 1 in"
                            "row/column spec \"{}\"".format(range_val, spec))

                # Store result into range_vals
                range_vals[i] = tmp


            # Range spec, create range
            if (len(range_vals) == 2 or len(range_vals) == 3):
                # Set increment range array
                if len(range_vals) == 3:
                    inc = abs(range_vals[2])
                else:
                    inc = 1

                # Array couting from range_vals[0] to range_vals[1]
                if (range_vals[0] <= range_vals[1]):
                    array += range(range_vals[0], range_vals[1]+1, inc)

                # Array couting down from range_vals[0] to range_vals[1]
                else:
                    array += range(range_vals[0], range_vals[1]-1,-inc)

            # Individual value
            elif len(range_vals) == 1:
                array += range_vals

            # Something is wrong
            else:
                raise ValueError('invalid number of range arguments "{}" '\
                    'in row/column spec "{}"'.format(list_val, spec))

        return array


class ImportXLSX (ImportSpreadsheet):
    def __init__(self, filename, sheet=None, entry=None):
        # Call init of inherited class
        super(self.__class__, self).__init__(filename, "xlsx", entry=entry)

        # Load workbook
        self.xl_wb = openpyxl.load_workbook(self.filename)

        # Set sheet name
        self.add_sheet(sheet)

    def add_sheet(self, sheet):
        """ Set sheet name for workbook

        Used to save the name of the sheet used in the workbook. This
        will be used when the workbook is read.

        Args:
            sheet: name of the sheet from the workbook
                (if None, the active sheet is used)
        """

        self.sheet = sheet

    def read(self):
        """ Get the data from xlsx file

        Read the data specified by the row and column array of
        the object and return the data in a 2-dimensional array.

        Return:
            2-dimensional array of entries containing data from spreadsheet

        """

        # Get workbook sheet
        if self.sheet is None:
            xl_sheet = self.xl_wb.active
        else:
            xl_sheet = self.xl_wb[self.sheet]

        # Read in the spreadsheet values and create PreprocessorEntry objects
        self.data = []
        for r in range(1,xl_sheet.max_row + 1):
            self.data.append([])
            for c in range(1,xl_sheet.max_column + 1):
                cell = xl_sheet.cell(row=r, column=c)

                tmp = PreprocessorEntry(self.entry.tag, parent=self.entry,
                        elem=self.entry.elem, value=cell.value)

                self.data[r-1].append(tmp)

        # Process the data and return it
        return self.get_data()

class ImportCSV(ImportSpreadsheet):
    """ This class extends ImportSpreadsheet to support csv files

    This class contains methods for reading the specified rows and columns
    of a csv file and returning the values.

    """
    def __init__(self, filename, entry=None):
        # Call init of inherited class
        super(self.__class__, self).__init__(filename, "csv", entry=entry)

    def read(self):
        """ Get the data from csv file

        Read the data specified by the row and column array of
        the object and return the data in a 2-dimensional array.

        Return:
            2-dimensional array of entries containing data from spreadsheet

        """

        file_array = []
        max_col = 0

        self.data = []

        # Read entire file into memory
        with open(self.filename) as csvfile:
            reader = csv.reader(csvfile)
            for r,row in enumerate(reader):
                self.data.append([])
                for val in row:
                    tmp = PreprocessorEntry(self.entry.tag, parent=self.entry,
                            elem=self.entry.elem, value=val)

                    self.data[r].append(tmp)

        # Process the data and return it
        return self.get_data()


class PresentationPreprocessor:
    """ Preprocess data for creating pptx presentation.

    This class is used to preprocess and store information from an input
    XML file for the pptx-creator project. One purpose of the preprocessor
    is to create a tree of PreprocessorEntry and PreprocessorValue objects
    from the XML file. Another purpose is to evaluate and substitute
    variables using the get, set, mod elements and the prepend, append
    attributes.

    """

    def __init__(self, source=None):
        self.tree = None

        if source:
            self.parse(source)

    def parse(self, source):
        """ Load external XML document into the preprocessor.

        Load XML document and parse for pptx presentation format. A tree is
        created to store all of the presentation information.

        Args:
            source: The source XML document

        """

        # Open xml input file
        self.source = source
        etree = ET.parse(source, parser=LineNumberingParser())

        # Initialize parsing structures
        self.tree = PreprocessorEntry("_root_")
        self.var_stack = VariableStack()

        # Create tree, starting at the root
        self._process_element(etree.getroot(), self.tree)

        # ensure tree contains root
        if (len(self.tree.data) <= 0):
            raise ValueError("no root entry found in \"{}\""\
                    "".format(self.source))

        # ensure only one root element
        if (len(self.tree.data) > 1):
            raise ValueError("presentation must be the only root element")

        # ensure tree root is the presentation element
        if (self.tree.data[0].tag != "presentation"):
            raise ValueError("presentation must be the root element\n{}"
                    "".format(self.error_info(self.tree.data[0])))

    def get_root(self):
        """ Get the root presentation entry of the tree

        Return the PreprocessorEntry object representing the presentation
        element of the input XML file.

        Return:
            root of input XML file

        """

        if self.tree is None:
            raise ValueError("must parse input file before getting root\n")

        return self.tree.data[0]

    def _process_element(self, elem, parent_entry, parent_elem=None):
        """ Process element and sub-elements.

        This function creates new PreprocessorEntry object for the specified
        element from the ElementTree, performs preprocessing tasks for
        the element, recursively calls this function for each of the element's
        attributes, recursively calls this function for each sub-element, and
        performs any postprocessing tasks for the element.

        Args:
            elem: ElementTree element to be processed
            parent_entry: Parent entry that is to hold the new PreprocessorEntry
                being created. The new PreprocessorEntry will be added to the
                end of the parent's data array.

        Kwargs:
            parent_elem: This is set when the current element is a temporary
                element created to represent an attribute of the parent element
                in the XML file.

        """

        # Create entry for element
        if not parent_elem is None:
            elem_entry = PreprocessorEntry(elem.tag,
                    parent=parent_entry, elem=parent_elem, is_attrib=True)
        else:
            elem_entry = PreprocessorEntry(elem.tag,
                    parent=parent_entry, elem=elem)

        # Perform tasks on element before calling children
        self._preprocess_element(elem_entry)

        # Process each attribute as a sub-element
        for item in elem.items():
            attrib = ET.Element(item[0])
            attrib.text = item[1]
            self._process_element(attrib, elem_entry, parent_elem=elem)

        # Add text as _value data under element
        elem_entry.add_text(elem.text)

        # Process subelements
        for child in elem:
            # Process child element
            self._process_element(child, elem_entry)

            # Get text after sub element
            elem_entry.add_text(child.tail)

        # Perform tasks on element after creating children
        self._postprocess_element(elem_entry)

    def _preprocess_element(self, elem_entry):
        """ Perform preprocessing for element.

        This is called by the _process_element function for each
        element before its sub-elements are processed.

        The following operations are performed by this function:
            - push new scope onto the variable stack

        Args:
            elem_entry: element entry in tree

        """

        # Append new dict onto the var_stack for the current scope variables
        self.var_stack.push()


    def _postprocess_element(self, elem_entry):
        """ Perform postprocessing for element.

        This is called by the _process_element function for each
        element after its sub-elements are processed.

        The following operations are performed by this function:
            - pop the current scope from the variable stack
            - process any get, mod, set elements
            - process any prepend, append elements (attributes)

        Args:
            elem_entry: element entry in tree

        """

        # Initialize values
        get_mod_set = None
        append_prepend = None

        # Pop current dict to return to scope of parent element
        self.var_stack.pop()

        # Determine variable action
        if (elem_entry.tag == "get"):
            get_mod_set = elem_entry.tag
        elif (elem_entry.tag == "mod"):
            get_mod_set = elem_entry.tag
        elif (elem_entry.tag == "set"):
            get_mod_set = elem_entry.tag
        elif (elem_entry.tag == "append"):
            append_prepend = elem_entry.tag
        elif (elem_entry.tag == "prepend"):
            append_prepend = elem_entry.tag

        # Get variable value and append/prepend
        if append_prepend:
            # Get variable name and lookup value
            var = elem_entry.get_values(join=True)
            try:
                val = self.var_stack.get(var)
            except ValueError as err:
                raise ValueError("var name \"{}\" does not exist; must set "\
                        "var before using \"{}\" element\n{}".format(var,
                            append_prepend, self.error_info(elem_entry)))

            # Create new PreprocessorValue object and add parent
            pp_val = PreprocessorValue(value=val)
            pp_val.parent = elem_entry.parent

            # Append or prepend data to parent's data array
            if append_prepend == "append":
                pp_val.parent.data.append(pp_val)
            elif append_prepend == "prepend":
                pp_val.parent.data.insert(0, pp_val)
            else:
                # Should never get here
                assert False

            # Delete entry from tree since we're evaluating it now
            elem_entry.delete()


        # Get, modify, or set the value of variable
        elif get_mod_set:
            # Find variable name
            var_array = elem_entry.get_values(tag="var",join=True)

            if (len(var_array) == 0):
                raise ValueError("cannot find name of var in \"{}\" "\
                        "element\n{}".format(get_mod_set, 
                            self.error_info(elem_entry)))

            if (len(var_array) > 1):
                raise ValueError("name of var specified {} times in \"{}\" "\
                        "element; can only provide one var name\n{}"\
                        "".format(len(var_array), get_mod_set,
                            self.error_info(elem_entry)))

            # Get variable and value
            var = var_array[0]
            val = elem_entry.get_values(join=True)

            # Ensure variable name is not empty string
            if (var == ""):
                raise ValueError("var name cannot be empty string in {} "\
                        "element\n{}".format(get_mod_set,
                            self.error_info(elem_entry)))

            # Get, mod, set value
            if (get_mod_set == "get"):
                val = self.var_stack.get(var)
                elem_entry.to_value(val)
            elif (get_mod_set == "mod"):
                val = self.var_stack.mod(var, val)
                elem_entry.delete()
            elif (get_mod_set == "set"):
                val = self.var_stack.set(var, val)
                elem_entry.delete()

        # Ensure presentation only has slides as children
        elif (elem_entry.tag == "presentation"):

            # Ensure there is at least one child
            if (len(elem_entry.data) == 0):
                raise ValueError("presentation must have at least one slide"\
                        "element\n{}"\
                        "".format(self.error_info(elem_entry)))

            # Ensure _root_ is parent
            if (elem_entry.parent.tag is None or
                    elem_entry.parent.tag != "_root_"):
                raise ValueError("presentation must be the root element\n{}"\
                        "".format(self.error_info(elem_entry)))

            # Ensure all children are slides
            for child in elem_entry.data:
                if (child.tag != "slide"):
                    raise ValueError("presentation can only have slide"\
                            "elements as children\n{}"\
                            "".format(self.error_info(child)))

        # Ensure slide's parent is presentation
        elif (elem_entry.tag == "slide"):

            # Ensure presentation is parent
            if (elem_entry.parent.tag is None or
                    elem_entry.parent.tag != "presentation"):
                raise ValueError("slide element must have presentation"\
                        "element as a parent\n{}"\
                        "".format(self.error_info(elem_entry)))

        # Ensure placeholder's parent is slide
        elif (elem_entry.tag == "placeholder"):

            # Ensure slide is parent
            if (elem_entry.parent.tag is None or
                    elem_entry.parent.tag != "slide"):
                raise ValueError("placeholder element can only have slide"\
                        "element as a parent\n{}"\
                        "".format(self.error_info(elem_entry)))

            # Convert placeholder tag to the name of the placeholder
            name_vals = elem_entry.get_values(tag="name", join=True)

            if (len(name_vals) > 1):
                raise ValueError("placeholder may only have one name attribute"\
                        "\n{}".format(self.error_info(elem_entry)))

            elem_entry.tag = name_vals[0]
            elem_entry.remove(tag="name")

    def error_info(self, elem_entry):
        return "\tFile \"{}\" {}".format(self.source, elem_entry.error_info())


class VariableStack:
    """ This is the variable stack class.

    This class is used to create and manage the variable stack for the input
    XML file. The variables that are defined in a scope are all stored in a
    single dictionary for that scope. A new sub-scope can be created by
    "push"-ing a new dictionary to the end of the variable stack. Any variable
    defined in this new scope will be available in all sub-scopes unless they
    are "set" in a sub-scope. In this case, the new value will be returned
    until this sub-scope is "pop"-ed and then the previous value will be
    returned. A "mod" can be used to modify the value of a variable in the
    current scope or any parent scopes, allow the value to be retained after
    the sub-scope is "pop"-ed.

    """

    def __init__(self):
        self.var_stack = []

    def push(self):
        """ Add new level to variable stack

        Add a new dictionary to the top of the variable stack to create a new
        variable scope.

        """

        # Add new dictionary to end of var_stack
        self.var_stack.append({})

    def pop(self):
        """ Remove level from variable stack

        Remove top dictionary from the variable stack to return to the previous
        variable scope.

        """

        # Remove dictionary from end of var_stack
        self.var_stack.pop()

    def set(self, var, val):
        """ Set variable value in the current scope

        Add new variable entry to the top dictionary of the variable stack,
        which will create a definition of the variable in the current scope

        Args:
            var: name of variable
            val: value of variable

        """

        # Set variable value to the dictionary at the end of var_stack
        self.var_stack[-1][var] = val

    def mod(self, var, val):
        """ Modify the value of the variable

        Find a previously defined variable in the dictionary closest to the
        top of the variable stack and modify the value.

        Args:
            var: name of variable
            val: new value of variable
        """

        # Change value of variable
        self.find_dict(var)[var] = val

    def get(self, var):
        """ Get the value of the variable

        Find a previously defined variable in the dictionary closest to the
        top of the variable stack and return the value.

        Args:
            var: name of variable

        Return:
            value of variable

        """

        # Get value of variable
        return self.find_dict(var)[var]

    def find_dict(self, var):
        """ Get the dictionary containing the specified variable

        Find a previously defined variable by searching through each dictionary
        from the top of the variable stack to the bottom. When the variable is
        found, the dictionary containing the variable is returned.

        Args:
            var: name of variable

        Return:
            dictionary containing variable

        """

        # Search through variable stack
        for var_dict in reversed(self.var_stack):
            if var in var_dict:
                return var_dict

        # Variable not found
        raise ValueError("var \"{}\" not found in variable stack".format(var))


class PreprocessorEntry:
    """ This is the preprocess entry class.

    The preprocess entry object is used by the PresentationPreprocessor
    class to hold elements and attributes in a standard data structure.

    An entry has a tag attribute and an entries attribute.
        - tag:  string containing attribute or element name from XML file
        - data: array of additional entries or values

    An object can be constructed with or without a value, as shown below
        entry = PreprocessorEntry("my_tag")
            Returns:
                entry.tag  == "my_tag"
                entry.data == []

        entry = PreprocessorEntry("my_tag",value="my_value")
            Returns:
                entry.tag  == "my_tag"
                entry.data[0].value == "my_value"

    """
    # Used to differentiate PreprocessorEntry and PreprocessorValue objects
    which = "entry"

    def __init__(self, tag, parent=None, value=None, elem=None, is_attrib=False):
        """ Initialize new PreprocessorEntry object.

        Create object with tag and empty data array. The capability is also
        provided to add a value associated with the enty, by setting value.
        When value is set, a new PreprocessorValue object is created with
        the value and added as the first entry of the ProprocessorEntry
        data array.

        The caller also has the option to associate an ElementTree element
        with this entry using the elem keyword argument. This can be done
        to provide clearer error messages and to aid in debugging. The
        is_attrib keyword argument can be used to indicate whether this entry
        is the element itself or an attribute of the element.

        Args:
            tag: The name of the tag for the entry which corresponds to the
                element tag or the attribute name.

        Kwargs:
            parent: The parent entry that is to hold the newly created entry.
                The new entry will be appended to the end of the data array
                of the parent entry.
            value: The value associated with the entry, typically this is
                used when adding an attribute entry so that both the name
                and value of the attribute can be added at once.
            elem: The ElementTree element corresponding to this entry. This
                is used for printing helpful error messages.
            is_attrib: True/False to indicate whether the entry is an
                attribute or an element in the original XML file.

        """

        # Initialize data
        self.tag = tag
        self.data = []
        self.elem = elem
        self.is_attrib = is_attrib

        # When value is specified create PreprocessorValue in data array
        if not value is None:
            self.add_value(value)

        # Add entry to end of data array for parent
        if parent:
            self.parent = parent
            parent.data.append(self)

    def __repr__(self):
        return "<{} '{}' at {}>"\
                "".format( self.__class__.__name__, self.tag, hex(id(self)))

    def __str__(self):
        ret  = "<{}: \'{}\'".format(self.which, self.tag)

        if self.elem is None:
            ret += ">"
        else:
            ret += ", line: {}>".format(self.elem._start_line_number)

        for pp in self.data:
            for line in str(pp).splitlines():
                ret += "\n   {}".format(line)

        return ret

    def add_value(self, value):
        """ Add value to entry.

        Add a new PreprocessorValue entry to the current PreprocessorEntry
        object's data array.

        Args:
            value: value to be added to entry data array

        """
        PreprocessorValue(parent=self, value=value)

    def add_text(self, text):
        """ Add text to entry.

        Perform text processing (e.g. remove leading and trailing spaces
        on string.  Then add a new PreprocessorValue entry to the current
        PreprocessorEntry object's data array with the value of the
        processed text.

        Args:
            text: The text string to be added to the tree

        """

        # Don't add when text is None
        if (text is None):
            return

        # TODO may need to make this more intelligent
        # Remove all leading and trailing whitespace characters
        text = text.strip()

        # Only add text entry when text is not empty string
        if (text != ""):
            self.add_value(text)

    def get_values(self, tag=None, join=False):
        """ Return value(s) of entry

        Get all values (PreprocessorValue object data) in the current
        entry's data array. When a tag is specified, all value entries
        associated with the tag are returned, as an array. When the
        join argument is True, all the values associated with
        a specific entry are joined together. Therefore, when
        tag is unspecified and join is True, all the values of the current
        entry are joined and returned as a single value. When join
        is False, an array of values is returned. When tag is specified and
        join is True, an array of the joined values for each entry
        associated with the tag are returned. When the tag is specified
        and join is False, an array of the arrays of values for each entry
        associated with the tag are returned.

        Kwargs:
            tag: name of tag to retreive values from
            join: indicates whether to join values

        Return:
            if tag=None and join=True:  single string of joined values
            if tag=None and join=False: array of values
            if tag specified and join=True:  array of joined values
            if tag specified and join=False: array of arrays of values

        """

        # Get values from entries matching tag in data array
        if tag:
            return [pp.get_values(join=join) for pp in self.data
                    if pp.which == "entry" and pp.tag == tag]

        # Get values in data array
        values = [str(pp.value) for pp in self.data if pp.which == "value"]

        if join:
            return ''.join(values)
        else:
            return values

    def delete(self):
        """ Delete entry from parent's data array.

        Effectively deletes the element from the whole entry tree by removing
        the reference from the parent's data array.
        """

        self.parent.data.remove(self)

    def remove(self, tag=None):
        """ Remove entry with tag from the current element's data array

        Deletes all entries with the specified tag from the whole entry tree
        by removing the reference from this entry's data array.

        Args:
            tag: name of entry to remove

        """

        if tag is None:
            return

        for child in self.data:
            if (child.which == "entry" and child.tag == tag):
                self.data.remove(child)

    def to_value(self, val):
        """ Convert the current entry to a value.

        Delete the current PreprocessorEntry from the parent's data array
        and replace it with a ProprocessorValue with the value specified
        by val.
        """

        # Get current index and replace it with new PreprocessorValue object
        idx = self.parent.data.index(self)
        pp_val = PreprocessorValue(value=val)
        self.parent.data[idx] = pp_val

        # Update parent pointer in new pp_val object
        pp_val.parent = self.parent

    def error_info(self):
        return "in element \"{}\", line {}"\
                "".format(self.parent.tag, self.elem._start_line_number)

class PreprocessorValue:
    """ This is the preprocess value class.

    The preprocess value object is used by the PresentationPreprocessor
    class to hold values of elements and data for an entry.

    A value object has a value attribute, which holds the text value of
    elements and attributes from the XML file.

    """
    # Used to differentiate PreprocessorEntry and PreprocessorValue objects
    which = "value"

    def __init__(self, parent=None, value=None):
        """ Initialize new PreprocessorValue object.

        Create object with parent and value. The value is stored in the
        value variable for the object. The parent specifies the parent
        PreprocessorEntry where the new new PreprocessorValue should be
        added. This function will append the new PreprocessorValue to the
        end of the parent's data array.
        When value is set, a new PreprocessorValue object is created with

        Kwargs:
            parent: The parent entry that is to hold the newly created value.
                The new value will be appended to the end of the data array
                of the parent entry.
            value: The value of the PreprocessorValue object.

        """

        if value:
            self.value = value

        # Add entry to end of data array for parent
        if parent:
            self.parent = parent
            parent.data.append(self)

    def __repr__(self):
        return "<{} '{}' at {}>"\
                "".format( self.__class__.__name__, self.value, hex(id(self)))

    def __str__(self):
        ret  = "<{}: \"{}\">".format(self.which, self.value)
        return ret


# From Duncan Harris on stack overflow: https://stackoverflow.com/questions/
#   6949395/is-there-a-way-to-get-a-line-number-from-an-elementtree-element
class LineNumberingParser(ET.XMLParser):
    # Python3
    def _start(self, *args, **kwargs):
        # Here we assume the default XML parser
        element = super(self.__class__, self)._start(*args, **kwargs)
        self._start_helper(element)
        return element

    # Python2
    def _start_list(self, *args, **kwargs):
        # Here we assume the default XML parser
        element = super(self.__class__, self)._start_list(*args, **kwargs)
        self._start_helper(element)
        return element

    # Function for both python3 and python2
    def _start_helper(self, element):
        # Copy element position attributes into output Elements
        element._start_line_number = self.parser.CurrentLineNumber
        element._start_column_number = self.parser.CurrentColumnNumber
        element._start_byte_index = self.parser.CurrentByteIndex


    def _end(self, *args, **kwargs):
        element = super(self.__class__, self)._end(*args, **kwargs)
        element._end_line_number = self.parser.CurrentLineNumber
        element._end_column_number = self.parser.CurrentColumnNumber
        element._end_byte_index = self.parser.CurrentByteIndex
        return element

# Run program
main()
