#!/usr/bin/env python

#
# File: pptx-annotate.py
# Author: amort
#
# Dependencies:
#   - python-pptx: Install with "pip install --user python-pptx"
#   - pathlib2: (version 3.x only) Install with "pip install --user pathlib2"
#
# Versions History:
#   Version 0.1 (9-12-2018)
#       Initial version.
#
# Description: This file takes a powerpoint template and annotates it with
#   the slide layout and placeholder indexes.
#
#   This code is based on the code from the python-pptx documentation.
#
#

from pptx import Presentation
import argparse

default_template = ""
default_filename = "template_annotated.pptx"

def main():
  global verbose

  # **** Setup argument parser ****
  parser = argparse.ArgumentParser(description="This program takes a pptx "
          "template file and annotates it with the slide layout and "
          "placeholder indexes.")
  parser.add_argument("-o", "--output", default=default_filename,
          help="pptx output file")
  parser.add_argument("-t", "--template", default=default_template,
          help="pptx template file")
  parser.add_argument("-v", "--verbose", action="store_true")
  args = parser.parse_args()


  # **** Get arguments ****
  verbose           = args.verbose
  filename_output   = args.output
  filename_template = args.template

  # **** Annotate power point template ****
  if verbose:
    print("INFO: Reading template " + filename_template + "...")

  # Create presentation using template
  if (filename_template != ""):
    prs = Presentation(filename_template)
  else:
    prs = Presentation()

  # Loop through layouts and see various elements
  for idx, _ in enumerate(prs.slide_layouts):
    slide = prs.slides.add_slide(prs.slide_layouts[idx])

    # Loop through placeholders
    for shape in slide.placeholders:
      if (shape.is_placeholder):
        phf = shape.placeholder_format

        try:
          shape.text = 'Layout: {}, Placeholder: {}, Type: {}, Shape: {}'.format(idx, phf.idx, phf.type, shape.name)
        except AttributeError:
          print('Type {} has no text attribute'.format(phf.type))

        print('Layout: {}, Placeholder: {}, Type: {}, Shape: {}'.format(idx, phf.idx, phf.type, shape.name))

  # Save file
  prs.save(filename_output)


# Run program
main()
