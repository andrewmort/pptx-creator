#!/usr/bin/env python

#
# File: pptx-annotate.py
# Author: amort
#
# Dependencies:
#   - python-pptx: Install with "pip install --user python-pptx"
#   - pathlib2: (version 3.x only) Install with "pip install --user pathlib2"
#
# Versions History:
#   Version 0.1 (9-12-2018)
#       Initial version.
#
# Description: This file takes a powerpoint template and annotates it with
#   the slide layout and placeholder indexes.
#
#   This code is based on the code from the python-pptx documentation.
#
#

from pptx import Presentation
import os
import argparse

# Python compatibility
try:
    input_ = raw_input
except NameError:
    input_ = input


def main():
    global verbose

    # **** Setup argument parser ****
    parser = argparse.ArgumentParser(description="This program takes a pptx "
            "template file and annotates it with the slide layout and "
            "placeholder indexes.")
    parser.add_argument("template", help="pptx template file")
    parser.add_argument("-o", "--output", default="",
            help="pptx output file [default: TEMPLATE_annotated.pptx]")
    parser.add_argument("-v", "--verbose", action="store_true")
    args = parser.parse_args()


    # **** Get arguments ****
    verbose           = args.verbose
    filename_output   = args.output
    filename_template = args.template

    # Create output filename from template filename
    if (filename_output == ""):
        paths = os.path.splitext(filename_template)
        filename_output = paths[0] + "_annotated" + paths[1]

    # Check if output file exists and prompt user to delete
    if (os.path.isfile(filename_output)):
        if (not query_yes_no('File {} exists. '
                'Overwrite?'.format(filename_output))):
            exit()


    # **** Annotate power point template ****
    if verbose:
        print("INFO: Reading template " + filename_template + "...")

    # Create presentation using template
    if (filename_template != ""):
        prs = Presentation(filename_template)
    else:
        prs = Presentation()

    # Loop through layouts and see various elements
    for idx, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[idx])

        # Loop through placeholders
        for shape in slide.placeholders:
            if (shape.is_placeholder):
                phf = shape.placeholder_format

                try:
                  shape.text = 'Layout idx: {}, Placeholder idx: {}'\
                      '\nType: {}, Shape: {}'.format(idx, phf.idx, phf.type, shape.name)
                except AttributeError:
                  print('Type {} has no text attribute'.format(phf.type))

                if verbose:
                    print('Layout idx: {}, Placeholder idx: {}, '
                        'Type: {}, Shape: {}'.format(idx, phf.idx, phf.type, shape.name))

    # Save file
    prs.save(filename_output)
    print('Created annotated file {}.'.format(filename_output))

# Function: query_yes_no
# Copied code from user Bryce Guinta from
#   https://stackoverflow.com/questions/3041986/
#       apt-command-line-interface-like-yes-no-input
def query_yes_no(question, default=True):
    """Ask a yes/no question via standard input and return the answer.

    If invalid input is given, the user will be asked until
    they acutally give valid input.

    Args:
        question(str):
            A question that is presented to the user.
        default(bool|None):
            The default value when enter is pressed with no value.
            When None, there is no default value and the query
            will loop.
    Returns:
        A bool indicating whether user has entered yes or no.

    Side Effects:
        Blocks program execution until valid input(y/n) is given.
    """
    yes_list = ["yes", "y"]
    no_list = ["no", "n"]

    default_dict = {  # default => prompt default string
        None: "[y/n]",
        True: "[Y/n]",
        False: "[y/N]",
    }

    default_str = default_dict[default]
    prompt_str = "%s %s " % (question, default_str)

    while True:
        choice = input_(prompt_str).lower()

        if not choice and default is not None:
            return default
        if choice in yes_list:
            return True
        if choice in no_list:
            return False

        notification_str = "Please respond with 'y' or 'n'"
        print(notification_str)


# Run program
main()
